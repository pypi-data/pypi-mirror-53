import argparse
import csv
import os
from pathlib import Path

import pptx
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Cm, Pt

from ..utils import load_summary, load_evaluations

TITLE_ONLY_SLIDE = 5
BLANK_SLIDE = 6
SLIDE_WIDTH = 9144000
SLIDE_HEIGHT = 5143500
TEMPLATE_NAME = str(Path(__file__).parent / "template.pptx")

FONT_NAME = "Calibri"
FONT_SIZE = Pt(14)


DEFAULT_CONTENT = [
    [
        {
            "type": "text",
            "params": {
                "left": 0.2,
                "top": 8.36,
                "width": 5.4,
                "height": 2.6,
                "text": "(1 is control, 2 is scizophrenia)",
            },
        },
        {"type": "image", "params": {"top": 3.46, "height": 5, "name": "time_mode"}},
        {
            "type": "image",
            "params": {"top": 9, "height": 5, "name": "factor_scatterplot"},
        },
    ],
    [{"type": "image", "params": {"height": 7.27, "name": "leverage"}}],
    [{"type": "image", "params": {"name": "factor_lineplot", "height": 5}}],
]


def generate_uniqueness_table(slide, uniqueness_information, num_rows=10):
    # Setup table
    row_height = Cm(0.78)

    num_cols = 4
    col_widths = [Cm(4)]*4

    top, left = Cm(2.5), Cm(1)
    tablewidth = sum(col_widths)
    tableheight = (num_rows + 1) * row_height

    tableshape = slide.shapes.add_table(
        num_rows + 1, num_cols, left, top, tablewidth, tableheight
    )
    table = tableshape.table

    # Insert table content
    column_names = {
        'Run Name': 'name',
        'SSE Difference': 'SSE_difference',
        'Fit Difference': 'fit_difference',
        'FMS': 'fms',
    }
    column_names = [
        'Run Name',
        'SSE Difference',
        'Fit Difference',
        'FMS'
    ]
    data_names = [
        'name',
        'SSE_difference',
        'fit_difference',
        'fms'
    ]
    format_options = [
        '',
        '.2e',
        '.2e',
        '.2f'
    ]

    for cell_title, cell in zip(column_names, table.rows[0].cells):
        frame = cell.text_frame
        frame.clear()
        run = frame.paragraphs[0].add_run()
        run.text = cell_title

        font = run.font
        font.name = FONT_NAME
        font.size = FONT_SIZE
        font.bold = True

    for i in range(num_rows):
        cellrow = table.rows[i + 1]
        for data_name, format_option, cell in zip(
                data_names, format_options, cellrow.cells
            ):
            cell_data = uniqueness_information[data_name][i]
 
            frame = cell.text_frame
            frame.clear()
            run = frame.paragraphs[0].add_run()
            run.text = f"{cell_data:{format_option}}"

            font = run.font
            font.name = FONT_NAME
            font.size = FONT_SIZE

    # Format table
    for col_width, column in zip(col_widths, table.columns):
        column.width = col_width


def generate_overview_table(slide, data_rows, column_names):
    # Setup table
    num_rows = len(data_rows) + 1
    row_height = Cm(0.78)

    num_cols = 5
    col_widths = list(map(Cm, [3.15, 2.95, 4, 4.15, 3.5]))

    top, left = 476130, 2373459
    tablewidth = sum(col_widths)
    tableheight = (num_rows + 1) * row_height

    tableshape = slide.shapes.add_table(
        num_rows, num_cols, left, top, tablewidth, tableheight
    )
    table = tableshape.table

    # Insert table content
    for cell_title, cell in zip(column_names, table.rows[0].cells):
        frame = cell.text_frame
        frame.clear()
        run = frame.paragraphs[0].add_run()
        run.text = cell_title

        font = run.font
        font.name = FONT_NAME
        font.size = FONT_SIZE
        font.bold = True

    for i, row_data in enumerate(data_rows):
        cellrow = table.rows[i + 1]
        for cell_data, cell in zip(row_data.values(), cellrow.cells):
            frame = cell.text_frame
            frame.clear()
            run = frame.paragraphs[0].add_run()
            run.text = cell_data

            font = run.font
            font.name = FONT_NAME
            font.size = FONT_SIZE

    # Format table
    for col_width, column in zip(col_widths, table.columns):
        column.width = col_width


def add_image(slide, experiment, image_params):
    height = Cm(image_params.get("height", 5))
    vis_path = Path(experiment) / "summaries" / "visualizations"
    try:
        name = str(next(vis_path.glob(image_params["name"] + "*")))
    except StopIteration:
        return

    image = slide.shapes.add_picture(name, 0, 0, height=height)

    if "left" not in image_params:
        left = (SLIDE_WIDTH - image.width) / 2
    else:
        left = Cm(image_params["left"])

    if "top" not in image_params:
        top = (SLIDE_HEIGHT - image.height) / 2
    else:
        top = Cm(image_params["top"])
    image.left = int(left)
    image.top = int(top)


def add_text(slide, experiment, text_params):
    text_box = slide.shapes.add_textbox(
        Cm(text_params["left"]),
        Cm(text_params["top"]),
        Cm(text_params["width"]),
        Cm(text_params["height"]),
    )
    summary = load_summary(experiment)

    text_frame = text_box.text_frame
    text_frame.text = text_params["text"].replace("{rank}", str(summary["model_rank"]))
    for paragraph in text_frame.paragraphs:
        paragraph.font.name = FONT_NAME
        paragraph.font.bold = True
        paragraph.font.size = Pt(text_params.get("size", 14))
        paragraph.alignment = PP_ALIGN.LEFT

    text_frame.vertical_anchor = MSO_ANCHOR.TOP


def generate_slide(slide, experiment, slide_content):
    for content in slide_content:
        if content["type"].lower() == "image":
            add_image(slide, experiment, content["params"])
        elif content["type"].lower() == "text":
            add_text(slide, experiment, content["params"])


def generate_presentation(
    pres, data_rows, column_names, experiment_folder, slides_content=None
):
    # Setup slide
    pres.slide_width = SLIDE_WIDTH
    pres.slide_height = SLIDE_HEIGHT

    if len(pres.slides) > 0:
        slide = pres.slides[0]
    else:
        slide = pres.slides.add_slide(pres.slide_layouts[BLANK_SLIDE])

    generate_overview_table(slide, data_rows, column_names)

    for experiment in sorted(experiment_folder.iterdir()):
        if (
            not experiment.is_dir()
            or not (experiment / "summaries/summary.json").is_file()
        ):
            continue

        summary = load_summary(experiment)
        model = summary["model_type"].replace("_", " ")

        for slide_content in slides_content:
            slide = pres.slides.add_slide(pres.slide_layouts[TITLE_ONLY_SLIDE])
            slide.shapes.title.text = (
                f'{model} model with {summary["model_rank"]} components'
            )
            for paragraph in slide.shapes.title.text_frame.paragraphs:
                paragraph.font.name = FONT_NAME
                paragraph.font.bold = True
                paragraph.font.size = Pt(18)
                paragraph.alignment = PP_ALIGN.LEFT

            slide.shapes.title.text_frame.vertical_anchor = MSO_ANCHOR.TOP
            generate_slide(slide, experiment, slide_content)
        
        evaluations = load_evaluations(experiment)
        if (
            "multi_run_evaluations" in evaluations and
            "Uniqueness" in evaluations["multi_run_evaluations"]
        ):
            slide = pres.slides.add_slide(pres.slide_layouts[TITLE_ONLY_SLIDE])
            slide.shapes.title.text = (
                f'{model} model with {summary["model_rank"]} components'
            )
            for paragraph in slide.shapes.title.text_frame.paragraphs:
                paragraph.font.name = FONT_NAME
                paragraph.font.bold = True
                paragraph.font.size = Pt(18)
                paragraph.alignment = PP_ALIGN.LEFT

            generate_uniqueness_table(slide, evaluations["multi_run_evaluations"]["Uniqueness"])


        # for i, image in enumerate(sorted((experiment/'summaries'/'visualizations').iterdir())):
        #    slide = pres.slides.add_slide(pres.slide_layouts[TITLE_ONLY_SLIDE])
        #    slide.shapes.title.text = f'{model} model with {summary["model_rank"]} components'
        #    for paragraph in slide.shapes.title.text_frame.paragraphs:
        #    	paragraph.font.name = FONT_NAME
        #    	paragraph.font.bold = True
        #    	paragraph.font.size=Pt(18)
        #    	paragraph.alignment=PP_ALIGN.LEFT

        #    slide.shapes.title.text_frame.vertical_anchor = MSO_ANCHOR.TOP

        #    image = slide.shapes.add_picture(str(image), Cm(i*2), Cm(i*2), height=Cm(5))
        #    image.left = int((SLIDE_WIDTH - image.width)/2)
        #    image.top = int((SLIDE_HEIGHT - image.height)/2)

    return pres


def create_ppt(
    parent_folder, csvpath="slide.csv", pptpath="summary.pptx", slides_content=None
):
    if slides_content is None:
        slides_content = DEFAULT_CONTENT

    if parent_folder is not None:
        csvpath = os.path.join(parent_folder, csvpath)
        pptpath = os.path.join(parent_folder, pptpath)

    with open(csvpath) as f:
        reader = csv.DictReader(f)
        data_rows = [row for row in reader]
        column_names = list(data_rows[0].keys())

    pres = pptx.Presentation(TEMPLATE_NAME)
    pres = generate_presentation(
        pres, data_rows, column_names, parent_folder, slides_content=slides_content
    )
    pres.save(pptpath)


if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("experiment_folder")
    parser.add_argument("csv_file")
    parser.add_argument("output_file")
    parser.add_argument("-p", "--parent_folder", default=None)
    parser.add_argument("-t", "--template", default="template.pptx")

    args = parser.parse_args()
    experiment_folder = Path(args.experiment_folder)
    csv_file = experiment_folder / args.csv_file
    output_file = args.output_file
    if args.parent_folder is not None:
        csv_file = os.path.join(args.parent_folder, csv_file)
        output_file = os.path.join(args.parent_folder, output_file)
    # Get slide data
    with open(csv_file) as f:
        reader = csv.DictReader(f)
        data_rows = [row for row in reader]
        column_names = list(data_rows[0].keys())

    pres = pptx.Presentation(args.template)
    pres = generate_presentation(pres, data_rows, column_names, experiment_folder)
    pres.save(output_file)
